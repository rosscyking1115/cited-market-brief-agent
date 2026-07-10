"""Editable PPTX morning-meeting pack (plan §3: editable exports, not chat output).

Slides: title (watermark + meta) -> one per section -> evidence ledger table(s) ->
needs-review -> disclaimer. All text lives in normal text frames so analysts can
edit every word. Art. 50 marking goes into core document properties.
"""

import io
from datetime import UTC, datetime

from app.exports.bundle import ExportBundle, strip_claim_refs

NAVY = (0x00, 0x47, 0x7B)
INK = (0x16, 0x16, 0x16)
MUTED = (0x54, 0x5B, 0x63)
AMBER = (0xD6, 0x55, 0x13)
GREEN = (0x24, 0x87, 0x4B)

_LEDGER_ROWS_PER_SLIDE = 7


def build_pptx(bundle: ExportBundle) -> bytes:
    from pptx import Presentation  # noqa: PLC0415
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=INK, mono=False):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        para = frame.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
        run.font.name = "Courier New" if mono else "Arial"
        return frame

    # --- Title slide ---
    slide = prs.slides.add_slide(blank)
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(0.7), Emu(91440 // 8), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*NAVY)
    bar.line.fill.background()
    add_text(
        slide,
        0.8,
        0.65,
        8,
        0.6,
        "Cited Market Brief Agent",
        size=24,
        bold=True,
        color=NAVY,
    )
    add_text(slide, 0.8, 1.7, 11.5, 1.0, "What changed since yesterday?", size=34, bold=True)
    add_text(
        slide,
        0.8,
        2.8,
        11.5,
        0.5,
        f"Watchlist: {bundle.watchlist} · generated "
        f"{bundle.generated_at.strftime('%Y-%m-%d %H:%M UTC')} · model {bundle.model} · "
        f"{len(bundle.supported_claims)} validated claims",
        size=12,
        color=MUTED,
    )
    add_text(
        slide,
        0.8,
        3.5,
        11.5,
        0.5,
        bundle.watermark,
        size=13,
        bold=True,
        color=GREEN if bundle.status == "approved" else AMBER,
    )
    if bundle.approval_line:
        add_text(slide, 0.8, 4.0, 11.5, 0.4, bundle.approval_line, size=11, color=MUTED)

    # --- Section slides ---
    for section in bundle.sections:
        slide = prs.slides.add_slide(blank)
        suffix = "  (analyst-edited)" if section.action == "edit" else ""
        add_text(slide, 0.6, 0.5, 12, 0.6, section.title + suffix, size=22, bold=True, color=NAVY)
        add_text(slide, 0.6, 1.4, 12.1, 5.0, strip_claim_refs(section.content), size=14)
        add_text(slide, 0.6, 6.9, 12, 0.4, bundle.watermark, size=9, color=AMBER)

    # --- Evidence ledger slides ---
    claims = bundle.supported_claims
    for start in range(0, len(claims), _LEDGER_ROWS_PER_SLIDE):
        chunk = claims[start : start + _LEDGER_ROWS_PER_SLIDE]
        slide = prs.slides.add_slide(blank)
        add_text(slide, 0.6, 0.4, 12, 0.5, "Evidence ledger", size=20, bold=True, color=NAVY)
        table_shape = slide.shapes.add_table(
            rows=len(chunk) + 1,
            cols=4,
            left=Inches(0.6),
            top=Inches(1.1),
            width=Inches(12.1),
            height=Inches(5.4),
        )
        table = table_shape.table
        for col, head in enumerate(("ID", "Claim", "Sources", "Validator")):
            cell = table.cell(0, col)
            cell.text = head
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
        table.columns[0].width = Inches(1.0)
        table.columns[1].width = Inches(6.1)
        table.columns[2].width = Inches(3.8)
        table.columns[3].width = Inches(1.2)
        for r, claim in enumerate(chunk, start=1):
            for c, value in enumerate((claim.cid, claim.text, "; ".join(claim.sources), "PASS")):
                cell = table.cell(r, c)
                cell.text = value
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(10)
        add_text(slide, 0.6, 6.9, 12, 0.4, bundle.watermark, size=9, color=AMBER)

    # --- Needs review ---
    if bundle.review_claims:
        slide = prs.slides.add_slide(blank)
        add_text(slide, 0.6, 0.4, 12, 0.5, "Needs review — not validated, do not cite", size=20, bold=True, color=AMBER)
        lines = "\n".join(f"⚑ {c.cid}  {c.text}  — {c.reason}" for c in bundle.review_claims)
        add_text(slide, 0.6, 1.2, 12.1, 5.4, lines, size=12)

    # --- Disclaimer ---
    slide = prs.slides.add_slide(blank)
    add_text(slide, 0.6, 0.5, 12, 0.5, "Disclosures", size=20, bold=True, color=NAVY)
    add_text(slide, 0.6, 1.3, 12.1, 4.5, bundle.disclaimer, size=12, color=MUTED)

    # Art. 50 machine-readable marking + provenance in document properties
    prs.core_properties.comments = (
        f"ai_generated=true; brief_id={bundle.brief_id}; model={bundle.model}; "
        f"prompt={bundle.prompt_version}; format=cited-market-brief-agent.pptx/v1"
    )
    prs.core_properties.created = datetime.now(UTC).replace(tzinfo=None)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
